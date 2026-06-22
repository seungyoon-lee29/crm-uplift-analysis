"""CRM uplift 프로젝트 발표 자료(.pptx) 생성.

차트 4개(docs/figures/*.png)를 임베드하고, 분석 내러티브를 13슬라이드로 구성한다.
실행: python scripts/build_ppt.py  →  report/CRM_uplift_발표.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
FIG = ROOT / "docs" / "figures"

# ── 디자인 시스템 (차트 색과 통일) ──
PRIMARY = RGBColor(0x4F, 0x46, 0xE5)
INK = RGBColor(0x0F, 0x17, 0x2A)
MUTED = RGBColor(0x64, 0x74, 0x8B)
DANGER = RGBColor(0xDC, 0x26, 0x26)
LIGHT = RGBColor(0xF8, 0xFA, 0xFC)
ACCENT_BG = RGBColor(0xEE, 0xF2, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
KR = "Apple SD Gothic Neo"

SW, SH = Inches(13.333), Inches(7.5)


def _set(run, *, size, bold=False, color=INK, font=KR):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def _box(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    return tb, tf


def _rect(slide, left, top, width, height, color):
    from pptx.enum.shapes import MSO_SHAPE
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def _bg(slide, color):
    _rect(slide, 0, 0, SW, SH, color)


def _footer(slide, idx):
    _, tf = _box(slide, Inches(0.5), Inches(7.0), Inches(9), Inches(0.4))
    r = tf.paragraphs[0].add_run()
    r.text = "CRM Uplift · 전체발송 vs ML 타겟팅 · Hillstrom 실데이터"
    _set(r, size=9, color=MUTED)
    _, tf2 = _box(slide, Inches(12.3), Inches(7.0), Inches(0.7), Inches(0.4))
    p = tf2.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    r2 = p.add_run(); r2.text = str(idx); _set(r2, size=9, color=MUTED)


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def title_bar(slide, title, kicker=None):
    _rect(slide, 0, 0, Inches(0.18), Inches(1.25), PRIMARY)
    _, tf = _box(slide, Inches(0.55), Inches(0.35), Inches(12.2), Inches(1.0))
    if kicker:
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = kicker.upper(); _set(r, size=12, bold=True, color=PRIMARY)
        p2 = tf.add_paragraph()
    else:
        p2 = tf.paragraphs[0]
    r2 = p2.add_run(); r2.text = title; _set(r2, size=28, bold=True, color=INK)


def bullets(slide, items, left, top, width, height, size=16, gap=10):
    _, tf = _box(slide, left, top, width, height)
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        if isinstance(it, tuple):
            text, color, bold = it
        else:
            text, color, bold = it, INK, False
        r = p.add_run(); r.text = "•  " + text
        _set(r, size=size, color=color, bold=bold)


def add_image(slide, path, left, top, width):
    return slide.shapes.add_picture(str(path), left, top, width=width)


# ────────────────────────────────────────────────────────────
def build() -> Path:
    prs = Presentation()
    prs.slide_width, prs.slide_height = SW, SH

    # 1. Title
    s = blank(prs); _bg(s, INK)
    _rect(s, 0, Inches(3.05), SW, Inches(0.06), PRIMARY)
    _, tf = _box(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(1.1))
    r = tf.paragraphs[0].add_run(); r.text = "CRM 푸시, 다 보낼까 거를까?"
    _set(r, size=44, bold=True, color=WHITE)
    _, tf = _box(s, Inches(0.9), Inches(3.25), Inches(11.5), Inches(1.0))
    r = tf.paragraphs[0].add_run()
    r.text = "uplift(증분)로 답하는 타겟팅 — 평균이 아니라 인과로, 결론을 조작하지 않고 정직하게"
    _set(r, size=18, color=RGBColor(0xCB, 0xD5, 0xE1))
    _, tf = _box(s, Inches(0.9), Inches(5.6), Inches(11.5), Inches(1.2))
    for i, t in enumerate(["가상 사전과제 · Hillstrom 6.4만 명 무작위배정 이메일 A/B",
                            "Python · scikit-uplift · 인과추론 · train/validation/test",
                            "github.com/seungyoon-lee29/crm-uplift-analysis"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = t
        _set(r, size=13, color=RGBColor(0x94, 0xA3, 0xB8))

    # 2. 과제 / R2
    s = blank(prs); _bg(s, WHITE); title_bar(s, "풀어야 할 문제 — \"다 보내 vs 거르기\"", "과제 정의")
    _rect(s, Inches(0.55), Inches(1.7), Inches(7.4), Inches(4.9), LIGHT)
    bullets(s, [
        ('"K은행 그로스팀은 매 캠페인 전체 고객에게 푸시를 쏜다.', INK, True),
        ('전체발송 vs ML로 거른 타겟발송, 무엇이 옳은가?"', INK, True),
        ("R2(원하는 상태): 같은 발송예산에서 순증분이익 최대화 + 채널수명 방어", PRIMARY, True),
        ("핵심방정식:", MUTED, False),
        ("순이익/명 = 증분구매 × AOV × 마진 − 발송비 − 수신거부율 × LTV", INK, False),
        ("분해: 퍼널(발송→visit→conversion→spend) × uplift 세그먼트 × 피로비", MUTED, False),
    ], Inches(0.9), Inches(2.0), Inches(6.8), Inches(4.4), size=16, gap=14)
    _rect(s, Inches(8.3), Inches(1.7), Inches(4.5), Inches(4.9), ACCENT_BG)
    _, tf = _box(s, Inches(8.6), Inches(2.0), Inches(4.0), Inches(4.3))
    r = tf.paragraphs[0].add_run(); r.text = "왜 어려운가"; _set(r, size=16, bold=True, color=PRIMARY)
    for t in ["오픈율·전환율(평균)은 의사결정 단위가 아니다",
              "필요한 건 '메시지가 행동을 바꿨는가' = 증분",
              "수신거부 1건은 모든 미래 캠페인을 잃는 비용",
              "→ 평균이 아닌 한계·증분으로 사고해야"]:
        p = tf.add_paragraph(); p.space_after = Pt(10)
        r = p.add_run(); r.text = "– " + t; _set(r, size=14, color=INK)
    _footer(s, 2)

    # 3. 접근 원칙
    s = blank(prs); _bg(s, WHITE); title_bar(s, "접근 — 화려한 모델이 아니라 정직한 분석 태도", "프레임")
    cards = [
        ("증분 사고", "treated − control(무작위배정)로 의사결정. 평균효과가 아니라 인과 증분."),
        ("결론 비조작", "합성 데이터로 인과를 심지 않음. 실 벤치마크(Hillstrom)로 검증."),
        ("검증 규율", "train/validation/test 분리. 정책 선택과 성과 평가를 절대 섞지 않음."),
        ("가정 투명", "마진·LTV·수신거부는 데이터 밖 → config로 분리하고 민감도로 다룸."),
    ]
    x = Inches(0.55)
    for i, (h, b) in enumerate(cards):
        cx = Inches(0.55 + i * 3.13)
        _rect(s, cx, Inches(2.1), Inches(2.95), Inches(3.6), LIGHT)
        _rect(s, cx, Inches(2.1), Inches(2.95), Inches(0.12), PRIMARY)
        _, tf = _box(s, cx + Inches(0.2), Inches(2.45), Inches(2.6), Inches(3.0))
        r = tf.paragraphs[0].add_run(); r.text = f"{i+1}. {h}"; _set(r, size=17, bold=True, color=INK)
        p = tf.add_paragraph(); p.space_before = Pt(8)
        r = p.add_run(); r.text = b; _set(r, size=13, color=MUTED)
    _footer(s, 3)

    # 4. 데이터
    s = blank(prs); _bg(s, WHITE); title_bar(s, "데이터 — Hillstrom 무작위배정 이메일 A/B", "데이터")
    bullets(s, [
        ("64,000 고객을 3군 무작위배정: Mens E-Mail / Womens E-Mail / No E-Mail", INK, True),
        ("treatment = 이메일 발송(1) vs 무발송(0) · 결과 = visit → conversion → spend (2주)", INK, False),
        ("무작위배정 → control(무발송)이 깨끗한 반사실(counterfactual) 제공 = 증분 측정 가능", PRIMARY, True),
        ("합성 안 씀: 인과를 손으로 심는 순환논리를 피하려 실 벤치마크 선택", MUTED, False),
        ("데이터에 없는 것(=가정): 마진율·발송비·수신거부율·LTV → config로 노출", MUTED, False),
        ("검증: 결측 0 · 퍼널 단조성(conversion⊆visit, spend⊆conversion) 통과", MUTED, False),
    ], Inches(0.7), Inches(2.0), Inches(12), Inches(4.5), size=17, gap=16)
    _footer(s, 4)

    # 5. 결과1 — ATE
    s = blank(prs); _bg(s, WHITE); title_bar(s, "결과 ① 이메일은 평균적으로 강하게 먹힌다", "발견 1/3")
    rows = [("결과", "treated", "control", "증분(ATE)", "상대"),
            ("visit", "16.7%", "10.6%", "+6.1%p", "+57%"),
            ("conversion", "1.07%", "0.57%", "+0.50%p", "+86%"),
            ("spend", "$1.25", "$0.65", "+$0.60", "+91%")]
    tbl = s.shapes.add_table(4, 5, Inches(0.7), Inches(2.1), Inches(8.0), Inches(2.6)).table
    for c in range(5):
        tbl.columns[c].width = Inches(1.6)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.fill.solid(); cell.fill.fore_color.rgb = PRIMARY if ri == 0 else (LIGHT if ri % 2 else WHITE)
            tf = cell.text_frame; tf.word_wrap = True
            r = tf.paragraphs[0].add_run(); r.text = val
            _set(r, size=14, bold=(ri == 0 or ci == 0), color=WHITE if ri == 0 else INK)
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    bullets(s, [
        ('→ "그냥 다 보내"가 그 자체로 틀린 건 아니다.', PRIMARY, True),
        ("결론을 미리 정하고 데이터를 심으면 그게 더 위험 (지적 정직성의 출발점)", MUTED, False),
        ("AOV(구매자 평균 spend) ≈ $116 · 전체 구매율 0.9%", MUTED, False),
    ], Inches(0.7), Inches(5.1), Inches(12), Inches(1.6), size=15, gap=8)
    _footer(s, 5)

    # 6. 모델 — Qini
    s = blank(prs); _bg(s, WHITE); title_bar(s, "모델 — uplift T-learner, 신호는 약하지만 랜덤보다 위", "모델")
    add_image(s, FIG / "qini_curve.png", Inches(7.0), Inches(1.7), Inches(6.0))
    bullets(s, [
        ("scikit-uplift T-learner(TwoModels, GradientBoosting)", INK, True),
        ("Qini AUC: validation 0.047 · test 0.046", PRIMARY, True),
        ("Hillstrom은 uplift 신호가 본래 약한 데이터(문헌 0.03~0.06)", MUTED, False),
        ("→ 랜덤보다 위지만 강하지 않다. 이 약함이 뒤의 핵심 교훈으로 이어진다", INK, False),
    ], Inches(0.7), Inches(2.2), Inches(6.0), Inches(4.0), size=16, gap=16)
    _footer(s, 6)

    # 7. 정책 함정 — cherry-pick
    s = blank(prs); _bg(s, WHITE); title_bar(s, "함정 — 분위 cherry-pick은 노이즈에 과적합한다", "발견 2/3")
    add_image(s, FIG / "decile_net_profit.png", Inches(7.0), Inches(1.9), Inches(6.0))
    bullets(s, [
        ("'한계순이익>0 분위만 발송'은 직관적이지만 위험", INK, True),
        ("빨강 = validation이 제외한 분위인데, 정작 test 최고 분위(1·3)", DANGER, True),
        ("분위별 증분구매 95% CI가 넓다 = 부호를 단정할 수 없다", MUTED, False),
        ("비연속 선택 = validation 노이즈 과적합 → 불안정한 정책", INK, False),
    ], Inches(0.7), Inches(2.3), Inches(6.0), Inches(4.0), size=16, gap=16)
    _footer(s, 7)

    # 8. 정책 개선 — top-k
    s = blank(prs); _bg(s, WHITE); title_bar(s, "개선 — 단조 top-k% 정책 (검증 규율 유지)", "발견 3/3")
    add_image(s, FIG / "topk_policy.png", Inches(7.0), Inches(1.9), Inches(6.0))
    bullets(s, [
        ("정책 = '상위 k%에 발송'(단조). k*는 validation, 평가는 test", INK, True),
        ("기본 이메일: k* = 100% → 전체발송과 동률($1,490)", PRIMARY, True),
        ("같은 test에서 cherry-pick 정책은 −$1,033 손실", DANGER, True),
        ("핵심 교훈: 약한 신호에선 '모델'보다 '정책의 단조성'이 손익을 가른다", INK, True),
    ], Inches(0.7), Inches(2.3), Inches(6.0), Inches(4.2), size=15, gap=14)
    _footer(s, 8)

    # 9. 민감도
    s = blank(prs); _bg(s, WHITE); title_bar(s, "진짜 답은 채널 피로비에 달려 있다", "의사결정")
    add_image(s, FIG / "fatigue_sensitivity.png", Inches(7.0), Inches(1.9), Inches(6.0))
    bullets(s, [
        ("수신거부 1건 비용(피로비)을 얼마로 보느냐가 답을 가른다", INK, True),
        ("전체발송이 순손실로 도는 임계 = 약 $0.17/통", DANGER, True),
        ("이메일($0.05)은 안전권 → 전체발송 합리적", INK, False),
        ("푸시·SMS(피로비 큼) → k*가 100%→5%로 하강, 타겟팅이 흑자 전환", PRIMARY, True),
    ], Inches(0.7), Inches(2.3), Inches(6.0), Inches(4.2), size=16, gap=14)
    _footer(s, 9)

    # 10. 결론 / 팀장 설득
    s = blank(prs); _bg(s, INK); _rect(s, 0, 0, Inches(0.18), Inches(1.25), PRIMARY)
    _, tf = _box(s, Inches(0.55), Inches(0.45), Inches(12), Inches(1.0))
    r = tf.paragraphs[0].add_run(); r.text = "결론 — 팀장 설득"; _set(r, size=28, bold=True, color=WHITE)
    msgs = [
        ("1. 이메일처럼 싼 채널은 검증해보면 답이 \"다 보내라\"로 나온다 (k*=100%).",),
        ("2. 단, 정책을 노이즈에 과적합시키면(분위 cherry-pick) 오히려 $1,033을 잃는다.",),
        ("3. 채널 피로비가 $0.17/통을 넘는 순간(푸시·SMS) 전체발송은 손실로 돌아서고,",),
        ("   그때 top-k 타겟팅의 가치가 살아난다 — 그 경계를 수치로 보여준다.",),
    ]
    _, tf = _box(s, Inches(0.7), Inches(2.0), Inches(12), Inches(3.2))
    for i, (t,) in enumerate(msgs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(16)
        r = p.add_run(); r.text = t
        _set(r, size=19, bold=(i in (0, 2)), color=WHITE if i in (0, 2) else RGBColor(0xCB, 0xD5, 0xE1))
    _rect(s, Inches(0.7), Inches(5.7), Inches(11.9), Inches(1.0), RGBColor(0x1E, 0x29, 0x3B))
    _, tf = _box(s, Inches(1.0), Inches(5.85), Inches(11.4), Inches(0.8))
    r = tf.paragraphs[0].add_run()
    r.text = "\"ML로 거르면 무조건 이긴다\"가 아니라, 가정·검증·정책 단조성을 함께 본다."
    _set(r, size=15, bold=True, color=PRIMARY)

    # 11. 정직한 한계
    s = blank(prs); _bg(s, WHITE); title_bar(s, "정직한 한계 — 결론이 선 천장을 명시한다", "한계")
    bullets(s, [
        ("수신거부·LTV가 데이터에 없다 → 채널피로비는 가정 → 단일 답이 아닌 민감도로 제시", INK, True),
        ("uplift 신호 약함(Qini 0.046) → 분위 증분 CI가 넓다 → '확정 sleeping-dog' 단정 회피", INK, True),
        ("2008년 단일 캠페인 → 시점·업종 일반화 제한, 방법론 데모로 해석", INK, True),
        ("실데이터가 있다면: 수신거부 로그·LTV로 피로비 실측 · 다회차로 발송빈도↔피로 인과", MUTED, False),
    ], Inches(0.7), Inches(2.1), Inches(12), Inches(4.2), size=17, gap=18)
    _footer(s, 11)

    # 12. 무엇을 보여주나 + 전략적 사고
    s = blank(prs); _bg(s, WHITE); title_bar(s, "이 프로젝트가 증명하는 것", "역량")
    pairs = [
        ("증분 사고", "평균이 아니라 treated−control 인과로 의사결정"),
        ("지적 정직성", "결론 미리 정하지 않음 · 합성으로 심지 않음 · ML 자동승리 거부"),
        ("검증 규율", "train/val/test 분리 · 정책 단조화 · 95% CI로 과대주장 차단"),
        ("의사결정 소유", "R2→Gap→3프레임 분해→민감도, 팀장 설득까지 수렴"),
    ]
    for i, (h, b) in enumerate(pairs):
        cy = Inches(2.1 + (i // 2) * 2.3)
        cx = Inches(0.7 + (i % 2) * 6.2)
        _rect(s, cx, cy, Inches(5.9), Inches(2.0), LIGHT)
        _rect(s, cx, cy, Inches(0.12), Inches(2.0), PRIMARY)
        _, tf = _box(s, cx + Inches(0.3), cy + Inches(0.25), Inches(5.3), Inches(1.6))
        r = tf.paragraphs[0].add_run(); r.text = h; _set(r, size=18, bold=True, color=PRIMARY)
        p = tf.add_paragraph(); p.space_before = Pt(6)
        r = p.add_run(); r.text = b; _set(r, size=14, color=INK)
    _footer(s, 12)

    # 13. Thank you
    s = blank(prs); _bg(s, INK)
    _rect(s, 0, Inches(3.5), SW, Inches(0.06), PRIMARY)
    _, tf = _box(s, Inches(0.9), Inches(2.5), Inches(11.5), Inches(1.0))
    r = tf.paragraphs[0].add_run(); r.text = "감사합니다"; _set(r, size=40, bold=True, color=WHITE)
    _, tf = _box(s, Inches(0.9), Inches(3.8), Inches(11.5), Inches(1.5))
    for i, t in enumerate(["코드·리포트: github.com/seungyoon-lee29/crm-uplift-analysis",
                            "공급측 짝 프로젝트: amazon-seller-entry-analysis (마켓플레이스 양면)",
                            "make all 로 전체 재현 가능"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = t; _set(r, size=15, color=RGBColor(0xCB, 0xD5, 0xE1))

    out = ROOT / "report" / "CRM_uplift_발표.pptx"
    out.parent.mkdir(exist_ok=True)
    prs.save(out)
    return out


if __name__ == "__main__":
    p = build()
    print(f"[saved] {p}")
